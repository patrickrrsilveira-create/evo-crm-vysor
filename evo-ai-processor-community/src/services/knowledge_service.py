import logging
import io
import base64
import tiktoken
from typing import List, Dict, Any, Optional
from pypdf import PdfReader
from openai import AsyncOpenAI
import traceback
import httpx
from bs4 import BeautifulSoup

from src.services.database_service import DatabaseService
from src.services.global_config_service import GlobalConfigService

logger = logging.getLogger(__name__)

class KnowledgeService:
    def __init__(self, db: DatabaseService):
        self.db = db
        self.config_service = GlobalConfigService()
        
        # Tokenizer for splitting
        try:
            self.tokenizer = tiktoken.get_encoding("cl100k_base")
        except Exception as e:
            logger.warning(f"Could not load cl100k_base tokenizer: {e}. Falling back to r50k_base")
            self.tokenizer = tiktoken.get_encoding("r50k_base")

    async def get_openai_client(self) -> AsyncOpenAI:
        """Initialize the OpenAI async client using the global configuration API key from database."""
        try:
            # The global config API hides sensitive keys, so we must query the database directly
            query = """
                SELECT name, serialized_value FROM installation_configs 
                WHERE name IN ('OPENAI_API_SECRET', 'OPENAI_API_KEY', 'OPENAI_API_URL')
            """
            results = await self.db.fetch_all(query)
            
            api_key = None
            base_url = None
            for row in results:
                if row.get('name') == 'OPENAI_API_SECRET':
                    api_key = row.get('serialized_value')
                elif row.get('name') == 'OPENAI_API_KEY' and not api_key:
                    api_key = row.get('serialized_value')
                elif row.get('name') == 'OPENAI_API_URL':
                    base_url = row.get('serialized_value')
            
            if api_key:
                import json
                try:
                    if isinstance(api_key, str):
                        # Ruby serialization usually wraps strings in quotes
                        if api_key.startswith('"') and api_key.endswith('"'):
                            api_key = json.loads(api_key)
                    if isinstance(base_url, str):
                        if base_url.startswith('"') and base_url.endswith('"'):
                            base_url = json.loads(base_url)
                except Exception:
                    pass
                
                if base_url and base_url.strip():
                    # Ensure base_url has a protocol
                    if not base_url.startswith('http://') and not base_url.startswith('https://'):
                        base_url = 'https://' + base_url
                    return AsyncOpenAI(api_key=api_key, base_url=base_url)
                return AsyncOpenAI(api_key=api_key)
                
        except Exception as e:
            logger.error(f"Error fetching OpenAI key from DB: {e}")
            
        raise ValueError("OpenAI API key is not configured in Global Configs.")

    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from a PDF file."""
        try:
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)
            text = ""
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\\n--- Page {i + 1} ---\\n" + page_text
            return text
        except Exception as e:
            logger.error(f"Error extracting PDF: {e}\\n{traceback.format_exc()}")
            raise ValueError(f"Failed to parse PDF file: {str(e)}")

    def extract_text_from_docx(self, file_content: bytes) -> str:
        """Extract text from a Word (.docx) file."""
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_content))
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            return "\n".join(paragraphs)
        except Exception as e:
            logger.error(f"Error extracting DOCX: {e}\\n{traceback.format_exc()}")
            raise ValueError(f"Failed to parse Word file: {str(e)}")

    def extract_text_from_xlsx(self, file_content: bytes) -> str:
        """Extract text from an Excel (.xlsx) file."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"=== Sheet: {sheet.title} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [str(c) for c in row if c is not None and str(c).strip()]
                    if row_vals:
                        lines.append(" | ".join(row_vals))
            return "\n".join(lines)
        except Exception as e:
            logger.error(f"Error extracting XLSX: {e}\\n{traceback.format_exc()}")
            raise ValueError(f"Failed to parse Excel file: {str(e)}")

    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from a PowerPoint (.pptx) file."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_content))
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"--- Slide {i + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n".join(lines)
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}\\n{traceback.format_exc()}")
            raise ValueError(f"Failed to parse PowerPoint file: {str(e)}")

    async def extract_text_from_image(self, file_content: bytes, content_type: str = "image/jpeg") -> str:
        """Extract text from an image using OpenAI Vision."""
        try:
            client = await self.get_openai_client()
            b64 = base64.b64encode(file_content).decode("utf-8")
            resp = await client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "Extract and transcribe ALL text visible in this image. If it is a diagram or chart, describe it in detail. Output only the extracted content, no extra commentary."},
                            {"type": "image_url", "image_url": {"url": f"data:{content_type};base64,{b64}"}},
                        ],
                    }
                ],
                max_tokens=4096,
            )
            return resp.choices[0].message.content or ""
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}\\n{traceback.format_exc()}")
            raise ValueError(f"Failed to extract text from image: {str(e)}")


    async def extract_text_from_url(self, url: str) -> str:
        """Fetch and extract visible text from a URL."""
        try:
            async with httpx.AsyncClient(timeout=10.0) as client:
                response = await client.get(url, follow_redirects=True)
                response.raise_for_status()
                
            soup = BeautifulSoup(response.text, "html.parser")
            
            # Remove scripts, styles, head, title, meta, etc.
            for element in soup(["script", "style", "head", "title", "meta", "[document]"]):
                element.extract()
                
            text = soup.get_text(separator=" ", strip=True)
            return text
        except Exception as e:
            logger.error(f"Error extracting URL {url}: {e}\\n{traceback.format_exc()}")
            raise ValueError(f"Failed to extract text from URL: {str(e)}")

    def chunk_text(self, text: str, max_tokens: int = 500, overlap_tokens: int = 50) -> List[str]:
        """Split text into chunks of maximum `max_tokens` with an overlap."""
        tokens = self.tokenizer.encode(text)
        chunks = []
        
        if not tokens:
            return chunks

        i = 0
        while i < len(tokens):
            chunk_tokens = tokens[i:i + max_tokens]
            chunk_text = self.tokenizer.decode(chunk_tokens)
            chunks.append(chunk_text)
            
            i += (max_tokens - overlap_tokens)
            
        return chunks

    async def generate_embeddings(self, chunks: List[str]) -> List[List[float]]:
        """Generate vector embeddings for a list of text chunks using OpenAI."""
        if not chunks:
            return []
            
        client = await self.get_openai_client()
        
        try:
            # Generate embeddings in batches if necessary, but OpenAI supports multiple inputs up to a certain limit.
            # We'll assume the chunks list isn't astronomically large. A robust implementation would batch this.
            response = await client.embeddings.create(
                model="text-embedding-ada-002",
                input=chunks
            )
            
            # Sort the data by index to ensure order is preserved
            sorted_data = sorted(response.data, key=lambda x: x.index)
            return [data.embedding for data in sorted_data]
            
        except Exception as e:
            logger.error(f"Error generating embeddings: {e}\\n{traceback.format_exc()}")
            raise ValueError(f"Failed to generate embeddings: {str(e)}")

    async def process_and_store_document(
        self, 
        knowledge_base_id: str, 
        document_id: str, 
        text_content: str
    ):
        """Process text, generate chunks, generate embeddings, and store in DB."""
        logger.info(f"Processing document {document_id} for KB {knowledge_base_id}")
        
        # 1. Chunk Text
        chunks = self.chunk_text(text_content, max_tokens=500, overlap_tokens=50)
        logger.info(f"Generated {len(chunks)} chunks for document {document_id}")
        
        if not chunks:
            logger.warning(f"No text extracted for document {document_id}")
            return
            
        # 2. Generate Embeddings
        embeddings = await self.generate_embeddings(chunks)
        
        # 3. Store chunks and embeddings in database
        from sqlalchemy import text
        try:
            for chunk, embedding in zip(chunks, embeddings):
                embedding_str = f"[{','.join(map(str, embedding))}]"
                
                query_sql = text("""
                    INSERT INTO knowledge_document_chunks (knowledge_document_id, content, embedding, created_at, updated_at) 
                    VALUES (:doc_id, :chunk, :embedding::vector, NOW(), NOW())
                """)
                
                self.db.execute(query_sql, {
                    "doc_id": document_id,
                    "chunk": chunk,
                    "embedding": embedding_str
                })
            
            self.db.commit()
            logger.info(f"Successfully stored {len(chunks)} embedded chunks for document {document_id}")
        except Exception as e:
            self.db.rollback()
            logger.error(f"Error storing document chunks in DB: {e}")
            raise

    async def extract_process_and_store_document(self, knowledge_base_id: str, document_id: str, content: bytes, filename_lower: str, content_type: str):
        from starlette.concurrency import run_in_threadpool
        try:
            text_content = ""
            if filename_lower.endswith(".pdf") or "pdf" in content_type:
                text_content = await run_in_threadpool(self.extract_text_from_pdf, content)

            elif filename_lower.endswith(".docx") or "wordprocessingml" in content_type or "msword" in content_type:
                text_content = await run_in_threadpool(self.extract_text_from_docx, content)

            elif filename_lower.endswith(".xlsx") or filename_lower.endswith(".xls") or "spreadsheetml" in content_type or "excel" in content_type:
                text_content = await run_in_threadpool(self.extract_text_from_xlsx, content)

            elif filename_lower.endswith(".pptx") or filename_lower.endswith(".ppt") or "presentationml" in content_type or "powerpoint" in content_type:
                text_content = await run_in_threadpool(self.extract_text_from_pptx, content)

            elif any(filename_lower.endswith(ext) for ext in (".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp")) \
                    or content_type.startswith("image/"):
                img_type = content_type if content_type.startswith("image/") else "image/jpeg"
                text_content = await self.extract_text_from_image(content, img_type)

            elif filename_lower.endswith(".txt") or "text/plain" in content_type:
                text_content = content.decode("utf-8")
                
            if not text_content or not text_content.strip():
                logger.warning(f"No text extracted for document {document_id}")
                return
                
            await self.process_and_store_document(knowledge_base_id, document_id, text_content)
        except Exception as e:
            logger.error(f"Error processing background document {document_id}: {e}\\n{traceback.format_exc()}")

    async def retrieve_knowledge(self, query: str, knowledge_base_ids: List[str], limit: int = 5) -> str:
        """Search for relevant chunks across given knowledge bases."""
        if not knowledge_base_ids:
            return ""

        logger.info(f"Retrieving knowledge for query: '{query}' in KBs: {knowledge_base_ids}")
        
        # Embed the query
        query_embeddings = await self.generate_embeddings([query])
        if not query_embeddings:
            return ""
            
        query_embedding = query_embeddings[0]
        embedding_str = f"[{','.join(map(str, query_embedding))}]"
        
        # The vector operator <-> is L2 distance, <=> is cosine distance, <#> is inner product.
        query_sql = """
            SELECT c.content, c.embedding <=> $1::vector AS distance
            FROM knowledge_document_chunks c
            JOIN knowledge_documents d ON c.knowledge_document_id = d.id
            WHERE d.knowledge_base_id = ANY($2::uuid[])
            ORDER BY c.embedding <=> $1::vector
            LIMIT $3
        """
        
        try:
            results = await self.db.fetch_all(query_sql, embedding_str, knowledge_base_ids, limit)
            
            if not results:
                return ""
                
            # Combine the content into a single string
            combined_text = "\n\n---\n\n".join([row['content'] for row in results])
            return combined_text
        except Exception as e:
            logger.error(f"Error retrieving knowledge from DB: {e}")
            return ""

    async def search_agent_knowledge(self, agent_bot_id: str, query: str, limit: int = 5) -> str:
        """Search for relevant chunks across all knowledge bases attached to an agent bot."""
        # Get all knowledge_base_ids for this agent_bot_id or ai_agent_id
        query_sql = """
            SELECT knowledge_base_id FROM knowledge_base_agent_bots WHERE agent_bot_id = $1::uuid
            UNION
            SELECT knowledge_base_id FROM knowledge_base_ai_agents WHERE ai_agent_id = $1::uuid
        """
        try:
            records = await self.db.fetch_all(query_sql, agent_bot_id)
            knowledge_base_ids = [str(record['knowledge_base_id']) for record in records]
            
            if not knowledge_base_ids:
                return ""
                
            return await self.retrieve_knowledge(query, knowledge_base_ids, limit)
        except Exception as e:
            logger.error(f"Error finding knowledge bases for agent: {e}")
            return ""

